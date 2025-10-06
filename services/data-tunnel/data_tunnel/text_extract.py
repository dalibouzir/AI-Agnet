from __future__ import annotations

import csv
import io
import json
import logging
import mimetypes
import subprocess
import tempfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Dict, List, Optional

import pytesseract
from PIL import Image

logger = logging.getLogger(__name__)

try:  # Optional dependency; used for PDFs with textual content.
    import pdfplumber
except ImportError:  # pragma: no cover - best effort fallback
    pdfplumber = None

try:  # Generic structured document parsing fallback.
    from unstructured.partition.auto import partition
except ImportError:  # pragma: no cover - best effort fallback
    partition = None

try:  # Document specific extractors.
    from docx import Document
except ImportError:  # pragma: no cover
    Document = None

try:
    from pptx import Presentation
except ImportError:  # pragma: no cover
    Presentation = None

try:
    from openpyxl import load_workbook
except ImportError:  # pragma: no cover
    load_workbook = None


@dataclass
class TextExtractionResult:
    text: str
    doc_type: str
    pages: List[str] = field(default_factory=list)
    tables: List[str] = field(default_factory=list)
    metadata: Dict[str, Any] = field(default_factory=dict)
    ocr_applied: bool = False
    ocr_confidence: float = 1.0


def _extract_pdf(data: bytes) -> TextExtractionResult:
    if not pdfplumber:
        logger.debug("pdfplumber not available; returning binary placeholder")
        return TextExtractionResult(text="", doc_type="pdf")

    pages: List[str] = []
    tables: List[str] = []
    with pdfplumber.open(io.BytesIO(data)) as pdf:
        for page in pdf.pages:
            text = page.extract_text() or ""
            pages.append(text.strip())
            table_texts: List[str] = []
            try:
                for table in page.extract_tables():
                    table_rows = ["\t".join(cell or "" for cell in row) for row in table]
                    table_texts.append("\n".join(table_rows))
            except Exception as exc:  # pragma: no cover - pdfplumber quirks
                logger.debug("Table extraction failed: %s", exc)
            tables.extend(table_texts)
    combined_text = "\n\n".join(filter(None, pages + tables))
    return TextExtractionResult(text=combined_text, doc_type="pdf", pages=pages, tables=tables)


def _extract_docx(data: bytes) -> TextExtractionResult:
    if not Document:
        return TextExtractionResult(text="", doc_type="docx")
    document = Document(io.BytesIO(data))
    paragraphs = [p.text.strip() for p in document.paragraphs if p.text.strip()]
    tables: List[str] = []
    for table in document.tables:
        rows = ["\t".join(cell.text.strip() for cell in row.cells) for row in table.rows]
        tables.append("\n".join(row for row in rows if row))
    text_body = "\n".join(paragraphs + tables)
    return TextExtractionResult(text=text_body, doc_type="docx", pages=paragraphs, tables=tables)


def _extract_txt(data: bytes) -> TextExtractionResult:
    try:
        text = data.decode("utf-8")
    except UnicodeDecodeError:
        text = data.decode("latin-1", errors="ignore")
    return TextExtractionResult(text=text, doc_type="txt", pages=text.splitlines())


def _extract_csv(data: bytes) -> TextExtractionResult:
    handle = io.StringIO(data.decode("utf-8", errors="ignore"))
    reader = csv.reader(handle)
    rows = [", ".join(cell.strip() for cell in row if cell.strip()) for row in reader]
    text = "\n".join(rows)
    return TextExtractionResult(text=text, doc_type="csv", pages=rows, tables=[text])


def _extract_pptx(data: bytes) -> TextExtractionResult:
    if not Presentation:
        return TextExtractionResult(text="", doc_type="pptx")
    presentation = Presentation(io.BytesIO(data))
    slides: List[str] = []
    for slide in presentation.slides:
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                slide_text.append(shape.text.strip())
        if slide_text:
            slides.append("\n".join(slide_text))
    combined = "\n\n".join(slides)
    return TextExtractionResult(text=combined, doc_type="pptx", pages=slides)


def _extract_xlsx(data: bytes) -> TextExtractionResult:
    if not load_workbook:
        return TextExtractionResult(text="", doc_type="xlsx")
    workbook = load_workbook(io.BytesIO(data), data_only=True, read_only=True)
    sheets_text: List[str] = []
    for sheet in workbook.worksheets:
        values = []
        for row in sheet.iter_rows(values_only=True):
            row_values = [str(cell).strip() for cell in row if cell is not None]
            if row_values:
                values.append(", ".join(row_values))
        if values:
            sheets_text.append(f"Sheet {sheet.title}\n" + "\n".join(values))
    text = "\n\n".join(sheets_text)
    return TextExtractionResult(text=text, doc_type="xlsx", pages=sheets_text, tables=sheets_text)


def _extract_image(data: bytes, languages: str) -> TextExtractionResult:
    with Image.open(io.BytesIO(data)) as image:
        text = pytesseract.image_to_string(image, lang=languages)
    confidence = 0.6 if text.strip() else 0.0
    return TextExtractionResult(text=text, doc_type="image", pages=[text], ocr_applied=True, ocr_confidence=confidence)


def _extract_with_unstructured(data: bytes, filename: str) -> TextExtractionResult:
    if not partition:
        return TextExtractionResult(text="", doc_type="binary")
    elements = partition(file=io.BytesIO(data), file_filename=filename)
    text_fragments = [element.text for element in elements if getattr(element, "text", None)]
    text = "\n".join(text_fragments)
    return TextExtractionResult(text=text, doc_type="unstructured", pages=text_fragments)


def _run_ocrmypdf(data: bytes, languages: str) -> Optional[TextExtractionResult]:
    if not pdfplumber:
        return None
    try:
        with tempfile.NamedTemporaryFile(suffix=".pdf") as src, tempfile.NamedTemporaryFile(suffix=".pdf") as dst:
            src.write(data)
            src.flush()
            cmd = [
                "ocrmypdf",
                "--force-ocr",
                "--quiet",
                "--language",
                languages,
                src.name,
                dst.name,
            ]
            completed = subprocess.run(cmd, capture_output=True, check=False)
            if completed.returncode != 0:
                logger.debug("ocrmypdf failed: %s", completed.stderr.decode("utf-8", errors="ignore"))
                return None
            dst.seek(0)
            with open(dst.name, "rb") as handle:
                return _extract_pdf(handle.read())
    except FileNotFoundError:
        logger.debug("ocrmypdf not installed; skipping PDF OCR")
    except Exception as exc:  # pragma: no cover - rare system level error
        logger.warning("OCR fallback for PDF failed: %s", exc)
    return None


def _extract_json(data: bytes) -> TextExtractionResult:
    try:
        parsed = json.loads(data.decode("utf-8", errors="ignore"))
    except json.JSONDecodeError:
        return TextExtractionResult(text="", doc_type="json")
    text = json.dumps(parsed, indent=2)
    return TextExtractionResult(text=text, doc_type="json", pages=text.splitlines())


def _extension_from_filename(filename: str | None) -> str:
    if not filename:
        return ""
    return Path(filename).suffix.lower().strip(".")


def extract_text(
    data: bytes,
    filename: str | None,
    mime: str | None,
    *,
    enable_ocr: bool = True,
    ocr_languages: str = "eng",
) -> TextExtractionResult:
    """Return a best-effort textual representation for arbitrary files."""

    if not data:
        return TextExtractionResult(text="", doc_type="binary")

    extension = _extension_from_filename(filename)
    mime = mime or mimetypes.guess_type(filename or "")[0] or "application/octet-stream"

    if mime in {"application/pdf"} or extension == "pdf":
        result = _extract_pdf(data)
        if enable_ocr and not result.text.strip():
            ocr_result = _run_ocrmypdf(data, ocr_languages)
            if ocr_result:
                ocr_result.doc_type = "pdf"
                ocr_result.ocr_applied = True
                ocr_result.ocr_confidence = 0.7 if ocr_result.text.strip() else 0.0
                return ocr_result
        return result

    if mime in {"application/vnd.openxmlformats-officedocument.wordprocessingml.document"} or extension in {"docx"}:
        return _extract_docx(data)

    if mime in {"text/csv", "application/csv"} or extension == "csv":
        return _extract_csv(data)

    if mime.startswith("text/") or extension in {"txt", "md", "rtf"}:
        return _extract_txt(data)

    if mime in {"application/vnd.openxmlformats-officedocument.presentationml.presentation"} or extension == "pptx":
        return _extract_pptx(data)

    if mime in {"application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"} or extension in {"xlsx", "xlsm"}:
        return _extract_xlsx(data)

    if mime.startswith("image/") or extension in {"png", "jpg", "jpeg", "tiff", "bmp"}:
        if enable_ocr:
            return _extract_image(data, ocr_languages)
        return TextExtractionResult(text="", doc_type="image")

    if mime in {"application/json"} or extension == "json":
        return _extract_json(data)

    fallback = _extract_with_unstructured(data, filename or "upload.bin")
    if fallback.text.strip():
        return fallback

    if enable_ocr:
        ocr_guess = _extract_image(data, ocr_languages)
        if ocr_guess.text.strip():
            ocr_guess.doc_type = "image"
            return ocr_guess

    return TextExtractionResult(text="", doc_type=extension or "binary")
